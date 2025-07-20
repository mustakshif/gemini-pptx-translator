#!/usr/bin/env python3
"""
PPTX Translate with Gemini API (Fixed Version)
Translate PowerPoint files using Google's Gemini API with improved text extraction.
"""

import os
import json
import hashlib
import argparse
import asyncio
import logging
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional
import time

import google.generativeai as genai
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table
from pptx.slide import Slide

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('translation_log.txt', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class ProgressTracker:
    """Track and display translation progress."""
    
    def __init__(self, total_items: int, description: str = "Translating"):
        self.total_items = total_items
        self.current_item = 0
        self.description = description
        self.start_time = time.time()
    
    def update(self, current: int, text: str = ""):
        """Update progress."""
        self.current_item = current
        elapsed = time.time() - self.start_time
        if self.current_item > 0:
            estimated_total = elapsed * self.total_items / self.current_item
            remaining = estimated_total - elapsed
            eta_str = f"ETA: {remaining:.0f}s"
        else:
            eta_str = "ETA: --"
        
        percentage = (self.current_item / self.total_items) * 100
        progress_bar = self._create_progress_bar(percentage)
        
        # Clear line and print progress
        sys.stdout.write(f"\r{self.description}: {progress_bar} {percentage:.1f}% ({self.current_item}/{self.total_items}) {eta_str}")
        if text:
            sys.stdout.write(f" | {text[:30]}{'...' if len(text) > 30 else ''}")
        sys.stdout.flush()
    
    def finish(self):
        """Finish progress display."""
        elapsed = time.time() - self.start_time
        sys.stdout.write(f"\r{self.description}: {'█' * 50} 100.0% ({self.total_items}/{self.total_items}) Completed in {elapsed:.1f}s\n")
        sys.stdout.flush()
    
    def _create_progress_bar(self, percentage: float) -> str:
        """Create a text-based progress bar."""
        bar_length = 50
        filled_length = int(bar_length * percentage / 100)
        bar = '█' * filled_length + '░' * (bar_length - filled_length)
        return bar

class GeminiTranslator:
    """Gemini API translator for PowerPoint presentations."""
    
    def __init__(self, api_key: str, model: str = "gemini-2.5-flash"):
        """
        Initialize the Gemini translator.
        
        Args:
            api_key: Gemini API key
            model: Gemini model to use (default: gemini-1.5-flash)
        """
        self.api_key = api_key
        self.model = model
        self.cache = {}
        self.cache_file = None
        
        # Configure Gemini
        genai.configure(api_key=api_key)
        
        # Available Gemini models (as of 2025)
        self.available_models = {
            "gemini-2.5-flash": "Latest Gemini 2.5 flash model (fastest and most efficient)",
            "gemini-2.5-pro": "Latest Gemini 2.5 pro model (most capable)",
            "gemini-2.0-flash-exp": "Previous experimental flash model",
            "gemini-2.0-flash": "Previous stable flash model",
            "gemini-2.0-pro": "Previous pro model",
            "gemini-1.5-flash": "Legacy flash model",
            "gemini-1.5-pro": "Legacy pro model",
            "gemini-1.5-flash-exp": "Legacy experimental flash",
            "gemini-1.5-pro-exp": "Legacy experimental pro"
        }
        
        if model not in self.available_models:
            logger.warning(f"Model {model} not in known models. Available: {list(self.available_models.keys())}")
    
    def get_model_info(self) -> Dict[str, str]:
        """Get information about available models."""
        return self.available_models
    
    def set_cache_file(self, filename: str, language: str):
        """Set cache file for this translation session."""
        file_hash = hashlib.md5(filename.encode()).hexdigest()[:8]
        self.cache_file = f"translation_cache_{Path(filename).stem}_{language}_{file_hash}.json"
        self._load_cache()
    
    def _load_cache(self):
        """Load existing translation cache."""
        if self.cache_file and os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r', encoding='utf-8') as f:
                    self.cache = json.load(f)
                logger.info(f"Loaded cache with {len(self.cache)} entries from {self.cache_file}")
            except Exception as e:
                logger.warning(f"Failed to load cache: {e}")
                self.cache = {}
    
    def _save_cache(self):
        """Save translation cache to file."""
        if self.cache_file:
            try:
                with open(self.cache_file, 'w', encoding='utf-8') as f:
                    json.dump(self.cache, f, ensure_ascii=False, indent=2)
                logger.info(f"Saved cache with {len(self.cache)} entries to {self.cache_file}")
            except Exception as e:
                logger.error(f"Failed to save cache: {e}")
    
    def _get_cache_key(self, text: str, target_language: str) -> str:
        """Generate cache key for text and target language."""
        return hashlib.md5(f"{text}_{target_language}".encode()).hexdigest()
    
    async def translate_text(self, text: str, target_language: str, context: str = "") -> str:
        """
        Translate text using Gemini API.
        
        Args:
            text: Text to translate
            target_language: Target language code (e.g., 'en', 'zh-CN', 'es')
            context: Context information for better translation
            
        Returns:
            Translated text
        """
        if not text.strip():
            return text
        
        cache_key = self._get_cache_key(text, target_language)
        if cache_key in self.cache:
            logger.debug(f"Cache hit for: {text[:50]}...")
            return self.cache[cache_key]
        
        try:
            # Prepare prompt for translation
            prompt = f"""
            You are a professional translator. Translate the following text to {target_language}.
            
            Context: {context}
            
            Text to translate: "{text}"
            
            Instructions:
            1. Maintain the original meaning and tone
            2. Preserve any formatting markers or special characters
            3. Keep the translation natural and fluent
            4. If the text contains placeholders or variables, keep them unchanged
            5. Return only the translated text, nothing else
            
            Translated text:
            """
            
            # Use Gemini API with timeout
            model = genai.GenerativeModel(self.model)
            
            try:
                response = await asyncio.wait_for(
                    asyncio.to_thread(model.generate_content, prompt),
                    timeout=30  # 30 second timeout
                )
                
                if response and response.text:
                    translated_text = response.text.strip()
                    self.cache[cache_key] = translated_text
                    logger.debug(f"Translated: {text[:50]}... -> {translated_text[:50]}...")
                    return translated_text
                else:
                    logger.warning(f"Empty response from Gemini API for: {text[:50]}...")
                    return text
                    
            except asyncio.TimeoutError:
                logger.error(f"Translation timeout for '{text[:50]}...'")
                return text
                
        except Exception as e:
            error_msg = str(e)
            if "quota" in error_msg.lower():
                logger.error(f"API quota exceeded for '{text[:50]}...'")
            elif "key" in error_msg.lower():
                logger.error(f"Invalid API key for '{text[:50]}...'")
            else:
                logger.error(f"Translation error for '{text[:50]}...': {e}")
            return text
    
    async def translate_batch(self, texts: List[str], target_language: str, context: str = "", verbose: bool = False) -> List[str]:
        """
        Translate a batch of texts.
        
        Args:
            texts: List of texts to translate
            target_language: Target language code
            context: Context information
            
        Returns:
            List of translated texts
        """
        if not texts:
            return []
        
        # Filter out empty texts and check cache
        non_empty_texts = [(i, text) for i, text in enumerate(texts) if text.strip()]
        results = [text for text in texts]  # Initialize with original texts
        
        if not non_empty_texts:
            return results
        
        total_texts = len(non_empty_texts)
        logger.info(f"Starting translation of {total_texts} texts to {target_language}")
        
        # Initialize progress tracker (only if not verbose)
        progress = None if verbose else ProgressTracker(total_texts, f"Translating to {target_language}")
        
        # Translate non-empty texts with progress
        for idx, (i, text) in enumerate(non_empty_texts, 1):
            try:
                if verbose:
                    logger.info(f"Translating {idx}/{total_texts}: {text[:50]}{'...' if len(text) > 50 else ''}")
                else:
                    # Update progress bar
                    progress.update(idx, text)
                
                # Translate text
                results[i] = await self.translate_text(text, target_language, context)
                
                # Small delay to avoid rate limiting
                await asyncio.sleep(0.1)
                
            except Exception as e:
                logger.error(f"Failed to translate text {idx}: {e}")
                results[i] = text  # Keep original text on error
        
        # Finish progress display
        if progress:
            progress.finish()
        logger.info(f"Translation completed: {total_texts} texts processed")
        return results

class PPTXProcessor:
    """Process PowerPoint presentations for translation."""
    
    def __init__(self, translator: GeminiTranslator):
        self.translator = translator
    
    def extract_text_from_slide(self, slide: Slide) -> List[Dict[str, Any]]:
        """
        Extract text from a slide with improved text merging.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of text elements with metadata
        """
        text_elements = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 按段落提取完整文本，而不是按run分割
                for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        text_elements.append({
                            'text': paragraph_text,
                            'shape': shape,
                            'paragraph': paragraph,
                            'paragraph_idx': paragraph_idx,
                            'type': 'text'
                        })
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        # 提取单元格的完整文本
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_elements.append({
                                'text': cell_text,
                                'shape': shape,
                                'table': table,
                                'cell': cell,
                                'row_idx': row_idx,
                                'col_idx': col_idx,
                                'type': 'table'
                            })
        
        return text_elements
    
    def extract_all_text(self, presentation: Presentation) -> List[Dict[str, Any]]:
        """
        Extract all text from the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            
        Returns:
            List of all text elements with metadata
        """
        all_text_elements = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_elements = self.extract_text_from_slide(slide)
            for element in slide_elements:
                element['slide_number'] = slide_num + 1
                all_text_elements.append(element)
        
        return all_text_elements
    
    def apply_translation_to_paragraph(self, paragraph, translated_text: str):
        """
        Apply translated text to a paragraph while preserving formatting.
        
        Args:
            paragraph: PowerPoint paragraph object
            translated_text: Translated text to apply
        """
        # 清除现有的runs
        for run in paragraph.runs[::-1]:  # 反向删除避免索引问题
            paragraph._p.remove(run._r)
        
        # 添加翻译后的文本
        if translated_text:
            run = paragraph.add_run()
            run.text = translated_text
    
    def apply_translation_to_cell(self, cell, translated_text: str):
        """
        Apply translated text to a table cell.
        
        Args:
            cell: PowerPoint table cell object
            translated_text: Translated text to apply
        """
        # 清除单元格内容
        cell.text = ""
        
        # 设置翻译后的文本
        if translated_text:
            cell.text = translated_text
    
    async def translate_presentation(self, input_file: str, target_language: str, output_file: str = None, verbose: bool = False) -> str:
        """
        Translate a PowerPoint presentation.
        
        Args:
            input_file: Path to input PowerPoint file
            target_language: Target language code
            output_file: Path to output file (optional)
            
        Returns:
            Path to the translated presentation
        """
        logger.info(f"Starting translation of {input_file} to {target_language}")
        
        # Set cache file
        self.translator.set_cache_file(input_file, target_language)
        
        # Load presentation
        try:
            presentation = Presentation(input_file)
            logger.info(f"Loaded presentation with {len(presentation.slides)} slides")
        except Exception as e:
            logger.error(f"Failed to load presentation {input_file}: {e}")
            raise
        
        # Extract all text
        text_elements = self.extract_all_text(presentation)
        logger.info(f"Extracted {len(text_elements)} text elements")
        
        if not text_elements:
            logger.warning("No text found in presentation")
            return input_file
        
        # Translate texts
        texts_to_translate = [element['text'] for element in text_elements]
        translated_texts = await self.translator.translate_batch(
            texts_to_translate, 
            target_language, 
            context="PowerPoint presentation content",
            verbose=verbose
        )
        
        # Apply translations to presentation
        for element, translated_text in zip(text_elements, translated_texts):
            if element['type'] == 'text':
                self.apply_translation_to_paragraph(element['paragraph'], translated_text)
            elif element['type'] == 'table':
                self.apply_translation_to_cell(element['cell'], translated_text)
        
        # Save translated presentation
        if output_file is None:
            input_path = Path(input_file)
            output_file = str(input_path.parent / f"{input_path.stem}_translated_{target_language}{input_path.suffix}")
        
        try:
            presentation.save(output_file)
            logger.info(f"Saved translated presentation to {output_file}")
        except Exception as e:
            logger.error(f"Failed to save translated presentation: {e}")
            raise
        
        # Save cache
        self.translator._save_cache()
        
        return output_file

async def main():
    """Main function to handle command line interface."""
    parser = argparse.ArgumentParser(
        description="Translate PowerPoint presentations using Gemini API (Fixed Version)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_translate_gemini_fixed.py presentation.pptx -l en
  python pptx_translate_gemini_fixed.py presentation.pptx -l zh-CN -o translated.pptx
  python pptx_translate_gemini_fixed.py presentation.pptx -l es -m gemini-1.5-pro
  python pptx_translate_gemini_fixed.py -l fr  # Translate all .pptx files in current directory
        """
    )
    
    parser.add_argument('input_file', nargs='?', help='Input PowerPoint file (.pptx)')
    parser.add_argument('-l', '--language', required=True, help='Target language code (e.g., en, zh-CN, es, fr)')
    parser.add_argument('-o', '--output', help='Output file path')
    parser.add_argument('-k', '--api_key', help='Gemini API key (or set GEMINI_API_KEY env var)')
    parser.add_argument('-m', '--model', default='gemini-2.5-flash', 
                       help='Gemini model to use (default: gemini-2.5-flash)')
    parser.add_argument('--list-models', action='store_true', help='List available models and exit')
    parser.add_argument('--profile', action='store_true', help='Enable profiling')
    parser.add_argument('--verbose', action='store_true', help='Show detailed translation progress')
    parser.add_argument('--quiet', action='store_true', help='Minimize output (only errors and final result)')
    
    args = parser.parse_args()
    
    # List models if requested
    if args.list_models:
        translator = GeminiTranslator("dummy_key")  # Just to get model info
        print("Available Gemini models:")
        for model, description in translator.get_model_info().items():
            print(f"  {model}: {description}")
        return
    
    # Get API key
    api_key = args.api_key or os.getenv('GEMINI_API_KEY')
    if not api_key:
        logger.error("Gemini API key not provided. Set GEMINI_API_KEY environment variable or use -k option.")
        return
    
    # Set logging level based on arguments
    if args.quiet:
        logging.getLogger().setLevel(logging.ERROR)
    elif args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Initialize translator
    translator = GeminiTranslator(api_key, args.model)
    processor = PPTXProcessor(translator)
    
    # Handle input files
    if args.input_file:
        input_files = [args.input_file]
    else:
        # Find all .pptx files in current directory
        input_files = list(Path('.').glob('*.pptx'))
        if not input_files:
            logger.error("No .pptx files found in current directory")
            return
    
    # Process files
    start_time = time.time()
    for input_file in input_files:
        try:
            output_file = await processor.translate_presentation(
                str(input_file), 
                args.language, 
                args.output,
                verbose=args.verbose
            )
            logger.info(f"Successfully translated {input_file} -> {output_file}")
        except Exception as e:
            logger.error(f"Failed to translate {input_file}: {e}")
    
    if args.profile:
        elapsed_time = time.time() - start_time
        logger.info(f"Total processing time: {elapsed_time:.2f} seconds")

if __name__ == "__main__":
    asyncio.run(main()) 